# make_ppt.py
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def set_dark_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(11, 15, 25) # Sleek Slate/Dark Blue

def add_slide_title(slide, text, subtitle_text=None):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.name = "Segoe UI"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(16, 185, 129) # Cyber Emerald Green

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(6, 182, 212) # Cyber Neon Cyan

def build_standard_slide(prs, title, subtitle, bullets, img_filename):
    blank_layout = prs.slide_layouts[6] # Blank slide layout
    slide = prs.slides.add_slide(blank_layout)
    set_dark_bg(slide)
    add_slide_title(slide, title, subtitle)

    # Left Column: Bullet list of features
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.2), Inches(5.0))
    tf = left_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "•  " + bullet
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(226, 232, 240) # Off-white
        p.space_after = Pt(10)
        p.line_spacing = 1.15

    # Right Column: Live Screenshot
    img_path = os.path.join(os.path.dirname(__file__), 'screenshots', img_filename)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(6.0), Inches(1.5), width=Inches(6.8), height=Inches(4.85))
    else:
        err_box = slide.shapes.add_textbox(Inches(6.0), Inches(1.5), Inches(6.8), Inches(4.85))
        tf_err = err_box.text_frame
        tf_err.word_wrap = True
        p_err = tf_err.paragraphs[0]
        p_err.text = f"[Live Screenshot Pending: {img_filename}]"
        p_err.font.name = "Segoe UI"
        p_err.font.size = Pt(16)
        p_err.font.bold = True
        p_err.font.color.rgb = RGBColor(239, 68, 68)

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 Standard widescreen
    prs.slide_height = Inches(7.5)

    # 1. Slide: Title Slide
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_dark_bg(slide)

    # Glow visual design block (textbox)
    accent_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.5))
    tf = accent_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "EARTHSCAPE CLIMATE AGENCY (ECA)"
    p.font.name = "Segoe UI"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(16, 185, 129) # Emerald Green
    p.space_after = Pt(6)

    p2 = tf.add_paragraph()
    p2.text = "Distributed Climate Change Analytics & Visual Platform"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(6, 182, 212) # Cyan
    p2.space_after = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "Aptech Sem-6 eProject Curriculum Deliverable\nTechnology Stack: React 19, TypeScript, Vite, Recharts, Hadoop HDFS Core, MongoDB Catalog"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(148, 163, 184) # Slate gray
    p3.space_after = Pt(20)

    p4 = tf.add_paragraph()
    p4.text = "Group Project Deliverable | Submission Date: May 21, 2026"
    p4.font.name = "Segoe UI"
    p4.font.size = Pt(11)
    p4.font.italic = True
    p4.font.color.rgb = RGBColor(110, 231, 183) # Light green

    # 1b. Slide: Project Team & Timeline
    slide_team = prs.slides.add_slide(blank_layout)
    set_dark_bg(slide_team)
    add_slide_title(slide_team, "PROJECT TEAM & TIMELINE", "Aptech Sem-6 eProject (May 09, 2026 - June 08, 2026)")

    # Left Column: Team Members list
    left_box = slide_team.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0))
    tf_team = left_box.text_frame
    tf_team.word_wrap = True
    tf_team.margin_left = tf_team.margin_top = tf_team.margin_right = tf_team.margin_bottom = 0

    p_t = tf_team.paragraphs[0]
    p_t.text = "DEVELOPMENT TEAM MEMBERS:"
    p_t.font.name = "Segoe UI"
    p_t.font.size = Pt(15)
    p_t.font.bold = True
    p_t.font.color.rgb = RGBColor(6, 182, 212) # Cyan
    p_t.space_after = Pt(14)

    team_members = [
        ("MUHAMMAD OWAIS KHAN", "Student1483137"),
        ("MUHAMMAD OWAIS KHAN", "Student1484254"),
        ("MUHAMMAD SHEHARYAR", "Student1485732"),
        ("MUHAMMAD UMAR", "Student1484275"),
        ("MUHAMMAD ROHAN ZAFAR", "Student1486748"),
    ]

    for name, roll in team_members:
        p_m = tf_team.add_paragraph()
        p_m.text = f"•  {name}  ({roll})"
        p_m.font.name = "Segoe UI"
        p_m.font.size = Pt(13)
        p_m.font.color.rgb = RGBColor(226, 232, 240)
        p_m.space_after = Pt(10)

    # Right Column: Deliverables & Timeline Info
    right_box = slide_team.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.0))
    tf_time = right_box.text_frame
    tf_time.word_wrap = True
    tf_time.margin_left = tf_time.margin_top = tf_time.margin_right = tf_time.margin_bottom = 0

    p_tl = tf_time.paragraphs[0]
    p_tl.text = "PROJECT SCHEDULE & MILESTONES:"
    p_tl.font.name = "Segoe UI"
    p_tl.font.size = Pt(15)
    p_tl.font.bold = True
    p_tl.font.color.rgb = RGBColor(16, 185, 129) # Emerald
    p_tl.space_after = Pt(14)

    milestones = [
        "Start Date: 09-May-2026",
        "End Date: 08-Jun-2026",
        "Required Deliverables: Source Code, Compiled Code, DFDs, Flowcharts, and eProject Report.",
        "Core Analytics: Dashboard with dynamic Pearson-r engine.",
        "Cluster View: Hadoop HDFS DataNodes and MapReduce scheduler simulation.",
        "Predictive Analysis: LSTM temperature scenario generator.",
        "Support Center: Inquiry ticketing logging with mock responder."
    ]

    for ms in milestones:
        p_ms = tf_time.add_paragraph()
        p_ms.text = f"▪  {ms}"
        p_ms.font.name = "Segoe UI"
        p_ms.font.size = Pt(12)
        p_ms.font.color.rgb = RGBColor(226, 232, 240)
        p_ms.space_after = Pt(8)

    # 2. Slide: Problem Statement
    bullets_problem = [
        "Traditional IT upskilling suffers from rigid timings and far-flung practical environments.",
        "Students find it difficult to apply theoretical distributed storage concepts (like HDFS blocks) into practice.",
        "ECA resolves this challenge by delivering a high-fidelity visual simulation portal mapping actual big data architectures.",
        "Visualizing data flows (Level 0 and Level 1 DFDs) directly correlates the relationship between carbon footprints and surface heat."
    ]
    # Standard slides with screenshots
    build_standard_slide(
        prs,
        title="1. Project Context & Objectives",
        subtitle="Upskilling IT Professionals Through Interactive Ladders",
        bullets=bullets_problem,
        img_filename="01_landing_page.png"
    )

    # 3. Slide: Landing Page
    bullets_landing = [
        "A premium visual cyber-dark landing page acts as the main public website entrance.",
        "Features a dynamic 'Live Sentinel Feed' ticker simulating planetary telemetry in real-time.",
        "Showcases the three technological pillars (HDFS Core, ML Studio, and MongoDB Catalog).",
        "Hosts the main entry point to launch the secure analyst research portal."
    ]
    build_standard_slide(
        prs,
        title="2. Public Frontend Website",
        subtitle="Aesthetically Rich Entrypoint & Sentinel Ticker",
        bullets=bullets_landing,
        img_filename="01_landing_page.png"
    )

    # 4. Slide: Login Portal
    bullets_login = [
        "Restricted area login requiring decryption access keys to log in.",
        "Supports role-based authorization parameters distinguishing 'Admin' from 'Analyst' profiles.",
        "Provides quick-fill automated test buttons for fast academic evaluation.",
        "Logs invalid entries and features an integrated back button to return to the public site."
    ]
    build_standard_slide(
        prs,
        title="3. Access Control & Login Gate",
        subtitle="Secured Role-Based Authentication System",
        bullets=bullets_login,
        img_filename="02_login_page.png"
    )

    # 5. Slide: Tableau Dashboard
    bullets_dash = [
        "Interactive dashboard displaying dual multi-variable timelines.",
        "Calculates the Pearson Correlation Coefficient (r) dynamically based on selected date ranges.",
        "Displays real-time mathematical interpretations to evaluate surface heat against industrial carbon footprints.",
        "Features aggregate summary cards (Avg Temp, Avg CO2, Deforestation) using high-fidelity datasets."
    ]
    build_standard_slide(
        prs,
        title="4. Analytical Tableau Dashboard",
        subtitle="Real-time Pearson-r Mathematics Correlation Engine",
        bullets=bullets_dash,
        img_filename="03_dashboard_view.png"
    )

    # 6. Slide: Hadoop HDFS
    bullets_hadoop = [
        "Simulates a multi-node Hadoop storage rack with live CPU, RAM, and Disk capacity indicators.",
        "Interactive MapReduce visualizer illustrating map, shuffle, and reduce steps in real time.",
        "Prints clean pseudocode console logs detailing internal cluster tasks.",
        "Reinforces student understanding of distributed replication and split execution architectures."
    ]
    build_standard_slide(
        prs,
        title="5. Hadoop Distributed Storage Monitor",
        subtitle="Live MapReduce Thread Scheduler & Gauges",
        bullets=bullets_hadoop,
        img_filename="04_hadoop_hdfs_view.png"
    )

    # 7. Slide: ML Studio
    bullets_ml = [
        "Features hyperparameter sliders (learning rate, epochs, split) to simulate model training.",
        "Dynamic optimization plot showing training vs validation Mean Squared Error (MSE) convergence.",
        "Displays calculated mathematical metrics: R-squared index, MAE, and RMSE scores.",
        "Plots temperature forecasts up to 2050 based on Sustainable Mitigation vs. Business-As-Usual scenarios."
    ]
    build_standard_slide(
        prs,
        title="6. Predictive Machine Learning Studio",
        subtitle="Hyperparameter Tuning & LSTM Temperature Forecasting",
        bullets=bullets_ml,
        img_filename="05_predictive_ml_studio.png"
    )

    # 8. Slide: Ingestion Portal
    bullets_ingestion = [
        "Geospatial ingestion portal accepting weather CSVs, satellite GeoTIFF grids, and IoT JSON telemetry.",
        "Features validation stepper outlining data cleansing, HDFS 128MB segmentation, and rack placement.",
        "Logs database transactions syncing metadata directly into MongoDB collections.",
        "Simulates rack awareness policies by placing replication blocks on distinct datanodes."
    ]
    build_standard_slide(
        prs,
        title="7. Distributed Ingestion Portal",
        subtitle="Multi-Format Ingest, HDFS Block segmenter, and MongoDB Sync",
        bullets=bullets_ingestion,
        img_filename="06_ingestion_portal.png"
    )

    # 9. Slide: IoT Stream Center
    bullets_iot = [
        "Real-time sensor telemetry center fetching continuous temperature and carbon updates.",
        "Simulates remote weather station anomalies and raises visual safety alerts.",
        "Admin-only overrides: allows administrator profiles to override telemetry limits dynamically.",
        "Ensures immediate visual responses to anomalies to mock real-time disaster alerts."
    ]
    build_standard_slide(
        prs,
        title="8. Real-time IoT Stream Center",
        subtitle="Continuous Telemetry Stream & Alarm Override Gateways",
        bullets=bullets_iot,
        img_filename="07_iot_sensor_stream.png"
    )

    # 10. Slide: Support Desk
    bullets_support = [
        "Help desk view with a complete collapsible FAQ accordion covering cluster operations.",
        "Enables students to submit support tickets, generating customized ticket IDs.",
        "Displays correspondence logs, team responses, and active inquiry states.",
        "Mocks responsive agency communications for operational troubleshooting."
    ]
    build_standard_slide(
        prs,
        title="9. Support Desk & Interactive FAQ",
        subtitle="Inquiry Ticketing System & Troubleshooting Accordions",
        bullets=bullets_support,
        img_filename="08_support_and_faq.png"
    )

    # 11. Slide: Docs & DFDs
    bullets_docs = [
        "Embedded project manual explaining system architecture and requirements.",
        "Contains complete visual DFD Level 0 (Context) and Level 1 Diagrams designed in CSS.",
        "Embeds a visual MapReduce processing flowchart detailing distributed data paths.",
        "Satisfies academic requirements by presenting complete engineering structures in the UI."
    ]
    build_standard_slide(
        prs,
        title="10. Integrated Documentation Hub",
        subtitle="Context & Level 1 DFDs, MapReduce Flowcharts, and Hardware specs",
        bullets=bullets_docs,
        img_filename="09_documentation_hub.png"
    )

    # 12. Slide: Conclusion
    slide = prs.slides.add_slide(blank_layout)
    set_dark_bg(slide)

    conclusion_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.0))
    tf_c = conclusion_box.text_frame
    tf_c.word_wrap = True

    p_c1 = tf_c.paragraphs[0]
    p_c1.text = "THANK YOU"
    p_c1.font.name = "Segoe UI"
    p_c1.font.size = Pt(40)
    p_c1.font.bold = True
    p_c1.font.color.rgb = RGBColor(16, 185, 129)
    p_c1.alignment = PP_ALIGN.CENTER
    p_c1.space_after = Pt(14)

    p_c2 = tf_c.add_paragraph()
    p_c2.text = "EarthScape Climate Agency (ECA) eProject is Ready for Presentation.\n\nAll source code, setup files, running guidelines, DFD diagrams, and PPT materials are organized in the workspace folder."
    p_c2.font.name = "Segoe UI"
    p_c2.font.size = Pt(16)
    p_c2.font.color.rgb = RGBColor(226, 232, 240)
    p_c2.alignment = PP_ALIGN.CENTER

    try:
        prs.save(os.path.join(os.path.dirname(__file__), 'presentation.pptx'))
        print("PowerPoint presentation generated: presentation.pptx")
    except PermissionError:
        alternative_path = os.path.join(os.path.dirname(__file__), 'presentation_updated.pptx')
        prs.save(alternative_path)
        print(f"Permission denied on presentation.pptx (file might be open). Saved as: {os.path.basename(alternative_path)}")

if __name__ == "__main__":
    main()
